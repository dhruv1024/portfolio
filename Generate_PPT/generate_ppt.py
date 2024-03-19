from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Set slide width and height
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define slide layout with title and content
title_content_slide_layout = prs.slide_layouts[1]

# Custom colors
blue_color = RGBColor(0, 102, 204)
white_color = RGBColor(255, 255, 255)

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(title_content_slide_layout)
background = slide_1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = blue_color
title = slide_1.shapes.title
title.text = "Salary Review Request"
title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.color.rgb = white_color

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(title_content_slide_layout)
title_2 = slide_2.shapes.title
title_2.text = "Introduction"
content_2 = slide_2.placeholders[1]
content_2.text = (
    "I am writing to express my concerns regarding the salary evaluation in this year's appraisal process. "
    "I believe that my contributions and value to the firm may not be fully recognized and compensated in line with industry standards. "
    "Therefore, I kindly request a review of my current remuneration to ensure it aligns with market benchmarks and reflects my performance."
)
for paragraph in content_2.text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.size = Pt(22)

# Slide 3: Achievements and Contributions
slide_3 = prs.slides.add_slide(title_content_slide_layout)
title_3 = slide_3.shapes.title
title_3.text = "Achievements and Contributions"
content_3 = slide_3.placeholders[1]
content_3.text = (
    "As someone who recently assumed a new leadership role, I have actively driven efficiency within my team through initiatives in client management, mentorship, problem-solving, and engagement. "
    "This has resulted in tangible improvements and enhanced team performance.\n\n"
    "I believe that my achievements, as outlined in the appraisal form, and my consistent weighted appraisal/performance ratings, ranging between 3.75 and 4.03, indicate the value I bring to the firm."
)
for paragraph in content_3.text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.size = Pt(22)

# Slide 4: Dedication and Commitment
slide_4 = prs.slides.add_slide(title_content_slide_layout)
title_4 = slide_4.shapes.title
title_4.text = "Dedication and Commitment"
content_4 = slide_4.placeholders[1]
content_4.text = (
    "I have consistently demonstrated my dedication and commitment over the past five years, maintaining excellent attendance records and diligently fulfilling my responsibilities. "
    "With over 30+ PLs accrued and only a few leaves availed, predominantly due to weekend shifts, I have consistently prioritized my responsibilities and contributed to the team's success.\n\n"
    "Additionally, I have actively sought opportunities to enhance our support workflow, keeping abreast of industry advancements and implementing improvements that have positively impacted our processes."
)
for paragraph in content_4.text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.size = Pt(22)

# Slide 5: Conclusion and Request
slide_5 = prs.slides.add_slide(title_content_slide_layout)
title_5 = slide_5.shapes.title
title_5.text = "Conclusion and Request"
content_5 = slide_5.placeholders[1]
content_5.text = (
    "In order to maintain motivation and commitment, it is crucial that my compensation is fair and competitive, aligned with industry benchmarks. "
    "I am confident that my contributions have had a positive impact on the firm, and I am open to engage in further discussions and negotiations to ensure that it reflects my performance, responsibilities, and industry standards.\n\n"
    "I apologize for the slight delay in raising this matter. I took the time to thoroughly evaluate my performance and carefully consider the overall context before addressing this issue.\n\n"
    "Thank you for your attention to this matter. I look forward to a positive resolution."
)
for paragraph in content_5.text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.size = Pt(22)

# Save the presentation
prs.save("Salary_Review_Request_Presentation.pptx")
