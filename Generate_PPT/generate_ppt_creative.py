from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Set slide width and height
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define slide layout with title and content
title_content_slide_layout = prs.slide_layouts[1]
content_only_slide_layout = prs.slide_layouts[6]

# Custom colors
blue_color = RGBColor(0, 102, 204)
white_color = RGBColor(255, 255, 255)

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(title_content_slide_layout)
background = slide_1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = blue_color
title = slide_1.shapes.title
title.text = "Elevating My Worth: Salary Review Presentation"
title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.color.rgb = white_color

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(content_only_slide_layout)
content_2 = slide_2.placeholders[-1]
content_2.text = "Hello Team,\n\nI am thrilled to be here today to discuss my contributions, achievements, and aspirations within our company. Let's embark on this journey together."
content_2.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
content_2.text_frame.paragraphs[0].font.size = Pt(24)

# Slide 3: Achievements
slide_3 = prs.slides.add_slide(content_only_slide_layout)
content_3 = slide_3.placeholders[0]
content_3.text = "Recent Achievements\n\n- Implemented innovative client management strategies\n- Led successful mentorship programs\n- Introduced problem-solving initiatives resulting in efficiency gains\n- Enhanced team engagement leading to improved performance"
content_3.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
content_3.text_frame.paragraphs[0].font.size = Pt(24)

# Slide 4: Dedication and Commitment
slide_4 = prs.slides.add_slide(content_only_slide_layout)
content_4 = slide_4.placeholders[0]
content_4.text = "Dedication and Commitment\n\n- 5 years of dedicated service\n- Excellent attendance record\n- 30+ PLs accrued, demonstrating my commitment\n- Actively contributed to process improvements and workflow enhancements"
content_4.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
content_4.text_frame.paragraphs[0].font.size = Pt(24)

# Slide 5: Request for Salary Review
slide_5 = prs.slides.add_slide(content_only_slide_layout)
content_5 = slide_5.placeholders[0]
content_5.text = "Salary Review Request\n\nI kindly request a comprehensive salary review. My dedication, achievements, and commitment showcase my value to the team. Let's align my compensation with my contributions and aspirations. Thank you for your consideration."
content_5.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
content_5.text_frame.paragraphs[0].font.size = Pt(24)

# Save the presentation
prs.save("Creative_Salary_Review_Presentation.pptx")
